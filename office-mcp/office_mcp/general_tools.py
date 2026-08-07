"""General (cross-format) tools (VAL-GEN-001..025).

This module exposes the four **unprefixed** tools described in
``architecture.md`` §3.1 and validated by ``validation-contract.md``
(Area: General):

* :func:`list_documents`     — folder scan, non-recursive,
  ``type`` ∈ {``word``, ``excel``, ``pptx``}.
* :func:`get_document_info`  — auto-detect format by extension and
  dispatch to the format-specific ``*_get_info`` tool.
* :func:`search_text`        — cross-format text search reporting
  ``paragraph:N`` / ``cell:<coord>`` / ``slide:N:shape:M`` locations
  with a small surrounding-text context window.
* :func:`convert_document`   — dispatch to the format-specific
  exporter (``exporters`` module) by source extension and
  ``target_format``.

Conventions
-----------
* All four tools follow the same path-resolution and error contract
  as the per-format tools in ``word_tools``/``excel_tools``/
  ``pptx_tools``: :func:`office_mcp.paths.resolve_path`, raise
  :class:`OfficeMCPError` with the appropriate error code, return a
  plain ``dict``.
* ``list_documents`` is **non-recursive** (only the immediate
  children of ``folder``).
* ``get_document_info`` returns a dict whose ``type`` field is
  ``"word"``, ``"excel"``, or ``"pptx"`` plus the underlying
  ``*_get_info`` fields and a ``size_bytes`` field.
* ``search_text`` always returns ``[]`` for no matches and never
  raises on missing tokens.
* ``convert_document`` requires a non-empty ``target_format`` from
  ``{"pdf", "html", "csv"}``; ``csv`` is ``xlsx``-only and raises
  :data:`ERR_UNSUPPORTED_FMT` for any other source.
"""

from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from office_mcp.errors import (
    ERR_FILE_NOT_FOUND,
    ERR_INVALID_PARAMS,
    ERR_UNSUPPORTED_FMT,
    OfficeMCPError,
)
from office_mcp.paths import resolve_path
from server import mcp


# ---------------------------------------------------------------------------
# Module-level constants
# ---------------------------------------------------------------------------


#: Map of file extension → canonical type name used by ``list_documents``
#: and ``get_document_info``.
_TYPE_BY_EXT: dict[str, str] = {
    ".docx": "word",
    ".xlsx": "excel",
    ".pptx": "pptx",
}

#: Valid ``target_format`` values for :func:`convert_document`.
_VALID_TARGET_FORMATS: frozenset[str] = frozenset({"pdf", "html", "csv"})

#: Default number of characters of context to keep on each side of a
#: ``search_text`` match when building the ``context`` field.
_CONTEXT_PAD: int = 30


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _resolve_base_folder(folder: str | Path | None) -> Path:
    """Return an absolute, canonical path for the scan / output directory.

    * When ``folder`` is ``None``, fall back to
      :func:`office_mcp.config.get_default_folder`.
    * Relative ``folder`` values are joined onto the default folder so
      that the same default-folder semantics as the per-format tools
      apply.
    """
    if folder is None:
        from office_mcp.config import get_default_folder

        return get_default_folder()
    return Path(folder).expanduser().resolve()


def _build_context(text: str, start: int, end: int) -> str:
    """Return a window of up to ``_CONTEXT_PAD`` chars on each side.

    The match itself (``text[start:end]``) is included verbatim, so
    callers can substring-search for the original token inside the
    returned context.
    """
    ctx_start = max(0, start - _CONTEXT_PAD)
    ctx_end = min(len(text), end + _CONTEXT_PAD)
    return text[ctx_start:ctx_end]


def _match_locations(
    text: str,
    query: str,
    case_sensitive: bool,
) -> list[tuple[int, int, str]]:
    """Find every (non-overlapping) occurrence of ``query`` in ``text``.

    Returns a list of ``(start, end, context)`` tuples. ``context`` is
    a slice of the **original** text (case preserved) so the agent
    sees exactly what the user wrote.
    """
    if not query:
        return []
    if case_sensitive:
        haystack = text
        needle = query
    else:
        haystack = text.lower()
        needle = query.lower()
    if not needle:
        return []
    results: list[tuple[int, int, str]] = []
    pos = 0
    while True:
        idx = haystack.find(needle, pos)
        if idx < 0:
            break
        end = idx + len(needle)
        results.append((idx, end, _build_context(text, idx, end)))
        # Advance past the match — non-overlapping
        pos = end
    return results


# ---------------------------------------------------------------------------
# 1. list_documents
# ---------------------------------------------------------------------------


@mcp.tool()
def list_documents(folder: str | None = None) -> list[dict[str, Any]]:
    """List all Office files in a folder (non-recursive).

    Scans the immediate children of ``folder`` and returns one entry
    per file with extension ``.docx``/``.xlsx``/``.pptx``. Non-Office
    files (and subdirectories) are filtered out. The scan does **not**
    descend into subdirectories.

    Args:
        folder: Folder to scan. When ``None`` (default), the default
            folder from :func:`office_mcp.config.get_default_folder`
            is used. Relative paths are joined with the default folder.

    Returns:
        A list of dicts, each with:

        * ``path`` — absolute, canonical path to the file.
        * ``name`` — file basename (e.g. ``"report.docx"``).
        * ``type`` — one of ``"word"``, ``"excel"``, ``"pptx"``.
        * ``size_bytes`` — file size on disk (positive int).
        * ``modified`` — ISO-8601 timestamp of last modification
          (``datetime.fromisoformat``-parseable, with timezone).

        An empty folder returns ``[]`` (``VAL-GEN-002``).

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` if ``folder`` exists
            but is not a directory.
    """
    base = _resolve_base_folder(folder)
    if not base.exists():
        # Treating a missing folder the same as an empty folder keeps
        # the call idempotent (no side effects, no error) and is
        # friendlier to agents running against a fresh environment.
        return []
    if not base.is_dir():
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"Not a directory: {base}",
            {"path": str(base)},
        )
    results: list[dict[str, Any]] = []
    for entry in base.iterdir():
        if not entry.is_file():
            # Skip subdirectories, symlinks to dirs, etc. (non-recursive)
            continue
        suffix = entry.suffix.lower()
        type_name = _TYPE_BY_EXT.get(suffix)
        if type_name is None:
            continue
        try:
            stat = entry.stat()
        except OSError:
            # Race: file was removed between iterdir() and stat(). Skip.
            continue
        mtime = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
        results.append(
            {
                "path": str(entry.resolve()),
                "name": entry.name,
                "type": type_name,
                "size_bytes": stat.st_size,
                "modified": mtime.isoformat(),
            }
        )
    # Stable ordering by basename so the result is deterministic.
    results.sort(key=lambda e: e["name"])
    return results


# ---------------------------------------------------------------------------
# 2. get_document_info
# ---------------------------------------------------------------------------


def _dispatch_get_info(
    p: Path,
    type_name: str,
) -> dict[str, Any]:
    """Call the format-specific ``*_get_info`` tool and merge the result.

    Local imports avoid the circular dependency: the per-format
    modules are imported eagerly by ``server.py`` but they don't
    import from ``general_tools``.
    """
    if type_name == "word":
        from office_mcp.word_tools import word_get_info

        info = word_get_info(path=str(p))
        info["type"] = "word"
        return info
    if type_name == "excel":
        from office_mcp.excel_tools import excel_get_info

        info = excel_get_info(path=str(p))
        info["type"] = "excel"
        # ``excel_get_info`` returns a ``sheets`` list; expose the
        # summary fields that the contract lists (sheet_count,
        # sheet_names) for convenience. ``dimensions`` is also a
        # contract field — it is the union of all sheets'
        # ``ws.dimensions`` (e.g. ``"A1:C3"``).
        info["sheet_count"] = len(info["sheets"])
        info["sheet_names"] = [s["name"] for s in info["sheets"]]
        # Derive a single ``dimensions`` string by joining per-sheet
        # values. When the workbook has a single sheet (the common
        # case) this is just that sheet's ``dimensions`` attribute.
        try:
            wb = load_workbook(str(p), read_only=True)
            per_sheet = [
                f"{s['name']}!{wb[s['name']].dimensions}"
                for s in info["sheets"]
            ]
            wb.close()
        except Exception:  # pragma: no cover - defensive
            per_sheet = []
        info["dimensions"] = (
            per_sheet[0] if len(per_sheet) == 1 else ",".join(per_sheet)
        )
        return info
    # type_name == "pptx"
    from office_mcp.pptx_tools import pptx_get_info

    info = pptx_get_info(path=str(p))
    info["type"] = "pptx"
    # The contract asks for ``dimensions_inches``; ``pptx_get_info``
    # already returns ``dimensions`` with the same shape, but the
    # contract name is explicit so we expose both.
    info["dimensions_inches"] = info["dimensions"]
    return info


@mcp.tool()
def get_document_info(
    path: str,
    folder: str | None = None,
) -> dict[str, Any]:
    """Return metadata for a single Office file.

    Detects the format by extension and dispatches to the
    format-specific ``*_get_info`` tool. The returned dict is a
    superset of the underlying tool's dict, augmented with ``type``
    (``"word"`` / ``"excel"`` / ``"pptx"``) and ``size_bytes``.

    Args:
        path: Path to an existing Office file (``abs`` or relative to
            ``folder``/default folder).
        folder: Optional base folder for relative paths.

    Returns:
        A dict with:

        * ``type`` — one of ``"word"``, ``"excel"``, ``"pptx"``.
        * ``size_bytes`` — file size in bytes.
        * Format-specific fields: for ``.docx`` →
          ``paragraphs``, ``sections``, ``tables``, ``images``,
          ``properties``; for ``.xlsx`` → ``sheets`` (list),
          ``sheet_count``, ``sheet_names``; for ``.pptx`` →
          ``slide_count``, ``layouts``, ``dimensions_inches`` (with
          ``width_inches`` / ``height_inches``).

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if the file is missing
            (``VAL-GEN-010``), ``ERR_UNSUPPORTED_FMT`` if the
            extension is not one of ``.docx`` / ``.xlsx`` / ``.pptx``.
    """
    try:
        p = resolve_path(path, folder)
    except ValueError as exc:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            str(exc),
            {"path": path},
        ) from exc
    if not p.exists():
        raise OfficeMCPError(
            ERR_FILE_NOT_FOUND,
            f"File not found: {p}",
            {"path": str(p)},
        )
    suffix = p.suffix.lower()
    type_name = _TYPE_BY_EXT.get(suffix)
    if type_name is None:
        raise OfficeMCPError(
            ERR_UNSUPPORTED_FMT,
            f"Unsupported format: {suffix!r}",
            {"path": str(p), "extension": suffix},
        )
    size_bytes = p.stat().st_size
    info = _dispatch_get_info(p, type_name)
    info["size_bytes"] = size_bytes
    return info


# ---------------------------------------------------------------------------
# 3. search_text
# ---------------------------------------------------------------------------


def _search_docx(
    p: Path,
    query: str,
    case_sensitive: bool,
) -> list[dict[str, Any]]:
    doc = Document(str(p))
    results: list[dict[str, Any]] = []
    for i, paragraph in enumerate(doc.paragraphs):
        for _start, _end, context in _match_locations(
            paragraph.text, query, case_sensitive
        ):
            results.append(
                {
                    "location": f"paragraph:{i}",
                    "context": context,
                }
            )
    return results


def _search_xlsx(
    p: Path,
    query: str,
    case_sensitive: bool,
) -> list[dict[str, Any]]:
    wb = load_workbook(str(p), read_only=True, data_only=True)
    results: list[dict[str, Any]] = []
    try:
        for sheet in wb.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value is None:
                        continue
                    text = str(cell.value)
                    for _start, _end, context in _match_locations(
                        text, query, case_sensitive
                    ):
                        results.append(
                            {
                                "location": f"cell:{cell.coordinate}",
                                "context": context,
                            }
                        )
    finally:
        wb.close()
    return results


def _search_pptx(
    p: Path,
    query: str,
    case_sensitive: bool,
) -> list[dict[str, Any]]:
    prs = Presentation(str(p))
    results: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape_idx, shape in enumerate(slide.shapes):
            # ``text_frame`` is the usual path; fall back to ``shape.text``
            # for shapes (like placeholders) that have a ``.text`` shortcut
            # but no ``text_frame`` (rare in practice).
            if shape.has_text_frame:
                text = shape.text_frame.text
            else:
                try:
                    text = shape.text  # type: ignore[attr-defined]
                except AttributeError:
                    text = ""
            for _start, _end, context in _match_locations(
                text, query, case_sensitive
            ):
                results.append(
                    {
                        "location": f"slide:{slide_idx}:shape:{shape_idx}",
                        "context": context,
                    }
                )
    return results


@mcp.tool()
def search_text(
    path: str,
    query: str,
    case_sensitive: bool = True,
    folder: str | None = None,
) -> list[dict[str, Any]]:
    """Search for a substring across the body of an Office file.

    The location format depends on the source file's extension:

    * ``.docx`` → ``"paragraph:N"`` (0-based body paragraph index)
    * ``.xlsx`` → ``"cell:<coord>"`` (e.g. ``"cell:B3"``)
    * ``.pptx`` → ``"slide:N:shape:M"`` (0-based slide + 0-based shape
      index in ``slide.shapes``)

    Each match entry has a ``location`` and a ``context`` field. The
    ``context`` is a window of up to ~30 characters on each side of
    the match in the original text (case preserved).

    Args:
        path: Path to the file to search.
        query: Substring to look for. Must be a non-empty string.
        case_sensitive: When ``True`` (default), the match is exact;
            when ``False``, matches are case-insensitive.
        folder: Optional base folder for relative paths.

    Returns:
        A list of match dicts. Empty when there are no matches
        (``VAL-GEN-016``).

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` if ``query`` is not a
            string or is empty, ``ERR_FILE_NOT_FOUND`` if the file is
            missing, ``ERR_UNSUPPORTED_FMT`` for non-Office files.
    """
    if not isinstance(query, str):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"query must be a str, got {type(query).__name__}",
            {"type": type(query).__name__},
        )
    if not query:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            "query must be a non-empty string",
            {"query": query},
        )
    try:
        p = resolve_path(path, folder)
    except ValueError as exc:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            str(exc),
            {"path": path},
        ) from exc
    if not p.exists():
        raise OfficeMCPError(
            ERR_FILE_NOT_FOUND,
            f"File not found: {p}",
            {"path": str(p)},
        )
    suffix = p.suffix.lower()
    if suffix == ".docx":
        return _search_docx(p, query, case_sensitive)
    if suffix == ".xlsx":
        return _search_xlsx(p, query, case_sensitive)
    if suffix == ".pptx":
        return _search_pptx(p, query, case_sensitive)
    raise OfficeMCPError(
        ERR_UNSUPPORTED_FMT,
        f"search not supported for {suffix!r}",
        {"path": str(p), "extension": suffix},
    )


# ---------------------------------------------------------------------------
# 4. convert_document
# ---------------------------------------------------------------------------


def _default_output_path(src: Path, target: str) -> Path:
    """Sibling-of-input default: ``src.parent / (src.stem + '.<ext>')``."""
    return src.parent / f"{src.stem}.{target}"


def _convert_one(src: Path, target: str, out_dir: Path) -> Path:
    """Dispatch to the right ``exporters`` call for ``src`` / ``target``.

    Raises :class:`OfficeMCPError` with ``ERR_UNSUPPORTED_FMT`` when
    the (source-format, target) pair is not supported (e.g.
    ``.docx`` → ``.csv``).
    """
    # Local import: ``exporters`` is intentionally not imported at
    # module top so that import-time cycles cannot occur if the
    # per-format modules ever grow a dependency on ``general_tools``.
    from office_mcp import exporters

    suffix = src.suffix.lower()
    if target == "csv":
        # CSV is xlsx-only; raise early with a clear error before any
        # subprocess is launched.
        if suffix != ".xlsx":
            raise OfficeMCPError(
                ERR_UNSUPPORTED_FMT,
                f"CSV export is xlsx-only, got {suffix!r}",
                {"extension": suffix, "target": target},
            )
        return exporters.export_to_csv(src, out_dir)
    # target in {"pdf", "html"}: routed through the (universal) PDF
    # / HTML exporters, which already branch on the source extension.
    if target == "pdf":
        return exporters.export_to_pdf(src, out_dir)
    if target == "html":
        return exporters.export_to_html(src, out_dir)
    # Defensive — the caller validates ``target`` already.
    raise OfficeMCPError(
        ERR_UNSUPPORTED_FMT,
        f"Unsupported target format: {target!r}",
        {"target": target},
    )


@mcp.tool()
def convert_document(
    path: str,
    target_format: str,
    output: str | None = None,
    folder: str | None = None,
) -> dict[str, str]:
    """Convert a ``.docx``/``.xlsx``/``.pptx`` file to ``pdf``/``html``/``csv``.

    Dispatches to the format-specific exporter based on the source
    file's extension and the requested ``target_format``. When
    ``output`` is omitted, the output is written next to the source
    with the new extension (``VAL-GEN-024``).

    Args:
        path: Path to an existing Office file.
        target_format: One of ``"pdf"``, ``"html"``, ``"csv"``. The
            match is case-insensitive and a leading dot is tolerated.
        output: Optional target path for the converted file. When
            ``None`` (default), a sibling of the input is used.
            Relative paths are resolved against ``folder`` (or the
            default folder).
        folder: Optional base folder for relative paths.

    Returns:
        ``{"output_path": "<absolute path>", "format": <target>}``
        where ``format`` is the lower-cased, dot-stripped target
        format (e.g. ``"pdf"``).

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for unknown
            ``target_format``, ``ERR_FILE_NOT_FOUND`` if the source
            is missing, ``ERR_UNSUPPORTED_FMT`` for an unsupported
            (source, target) combination (e.g. ``.docx`` → ``.csv``,
            or any non-Office source), or any error propagated from
            :mod:`office_mcp.exporters` (``ERR_LIBREOFFICE_MISSING``,
            ``ERR_EXPORT_FAILED``).
    """
    if not isinstance(target_format, str):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"target_format must be a str, got {type(target_format).__name__}",
            {"type": type(target_format).__name__},
        )
    target = target_format.strip().lower().lstrip(".")
    if target not in _VALID_TARGET_FORMATS:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"Unknown target_format: {target_format!r}",
            {"target_format": target_format},
        )
    try:
        src = resolve_path(path, folder)
    except ValueError as exc:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            str(exc),
            {"path": path},
        ) from exc
    if not src.exists():
        raise OfficeMCPError(
            ERR_FILE_NOT_FOUND,
            f"File not found: {src}",
            {"path": str(src)},
        )
    suffix = src.suffix.lower()
    if suffix not in _TYPE_BY_EXT:
        raise OfficeMCPError(
            ERR_UNSUPPORTED_FMT,
            f"Unsupported source format: {suffix!r}",
            {"path": str(src), "extension": suffix},
        )
    # Resolve the output path. The "no output" default is sibling-of-input.
    if output is None:
        out_path = _default_output_path(src, target)
    else:
        out_path = resolve_path(output, folder)
    out_dir = out_path.parent
    produced = _convert_one(src, target, out_dir)
    # ``exporters`` may name the output ``<stem>.<ext>``; honour the
    # caller's requested filename by renaming if it differs.
    if produced != out_path:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        produced.replace(out_path)
        final = out_path
    else:
        final = produced
    return {"output_path": str(final), "format": target}
