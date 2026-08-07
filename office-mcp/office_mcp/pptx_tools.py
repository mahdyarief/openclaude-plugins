"""PowerPoint presentation tools for office-mcp.

This module exposes the 14 ``pptx_*`` tools described in
``architecture.md`` §3.4 (PowerPoint tools §1-14) and validated by
``validation-contract.md`` (Area: PowerPoint §1-14,
``VAL-PPTX-001``..``VAL-PPTX-065``).

The 7 core tools (§1-7, ``VAL-PPTX-001``..``-030``) provide the basic
CRUD operations on a presentation; the 7 content / export tools
(§8-14, ``VAL-PPTX-031``..``-065``) cover shape insertion
(text box / image / shape / table / chart) and PDF / HTML export.

Conventions
-----------
* Every tool resolves the user-supplied ``path`` (and optional
  ``folder``) through :func:`office_mcp.paths.resolve_path`.
* Every tool validates its inputs and raises :class:`OfficeMCPError`
  with the appropriate error code (``ERR_FILE_NOT_FOUND``,
  ``ERR_INVALID_PARAMS``, ``ERR_UNSUPPORTED_FMT``).
* Every tool opens the file with :class:`pptx.Presentation`,
  performs the operation, saves, and returns a plain ``dict`` (which
  FastMCP serializes as the ``structuredContent`` of the JSON-RPC
  response).
* ``pptx_create_presentation`` refuses to overwrite an existing file
  (``VAL-PPTX-003``); ``pptx_delete_slide`` refuses to delete the
  last remaining slide (``VAL-PPTX-024``); ``pptx_reorder_slides``
  is a *move* (not a swap — ``VAL-PPTX-027``).

Slide indexing
--------------
All slide indices are 0-based, matching :class:`pptx.Presentation`'s
``slides`` iteration. ``pptx_add_slide`` always appends to the end.
``pptx_delete_slide`` shifts later indices down. ``pptx_reorder_slides``
removes the slide at ``from_index`` first, then inserts it at
``to_index`` (so e.g. moving index 0 to index 2 in a 3-slide deck
produces the order ``[B, C, A]``, not the swap result ``[C, B, A]``).

Direct XML manipulation
-----------------------
:mod:`python-pptx` does not expose a built-in API to delete a slide
or to reorder slides. Both operations are implemented by
manipulating the ``sldIdLst`` XML element directly: the
:pyattr:`Presentation.slides._sldIdLst` attribute is a CT_SlideIdList
element whose children are ``sldId`` entries, one per slide. To
delete we remove the ``sldId`` element and drop its relationship
via :py:meth:`Part.drop_rel`; to reorder we move the ``sldId``
element to a new position in the list.

Inches, EMU, and DPI
--------------------
PowerPoint stores positions and sizes in English Metric Units (EMU).
1 inch = 914 400 EMU. We accept geometry arguments in inches and
convert via :class:`pptx.util.Inches`. For image natural size we
use 96 DPI (the standard for Office) — see
:func:`_picture_natural_size`.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Emu, Inches, Pt

from office_mcp.errors import (
    ERR_FILE_NOT_FOUND,
    ERR_INVALID_PARAMS,
    ERR_UNSUPPORTED_FMT,
    OfficeMCPError,
)
from office_mcp.paths import resolve_path, save_with_lock_check
from server import mcp


__all__ = [
    # Core (§1-7)
    "pptx_create_presentation",
    "pptx_get_info",
    "pptx_list_slides",
    "pptx_read_slide",
    "pptx_add_slide",
    "pptx_delete_slide",
    "pptx_reorder_slides",
    # Content (§8-12)
    "pptx_add_text_box",
    "pptx_add_image",
    "pptx_add_shape",
    "pptx_add_table",
    "pptx_add_chart",
    # Export (§13-14)
    "pptx_export_pdf",
    "pptx_export_html",
]


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _resolve_pptx(path: str, folder: str | None) -> Path:
    """Resolve ``path`` to an absolute ``Path`` and verify the extension.

    Wraps :func:`office_mcp.paths.resolve_path` so that empty /
    whitespace-only paths become :class:`OfficeMCPError`
    (``ERR_INVALID_PARAMS``) instead of a bare ``ValueError`, and
    enforces the ``.pptx`` extension up front so non-pptx files
    fail fast.
    """
    try:
        p = resolve_path(path, folder)
    except ValueError as exc:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            str(exc),
            {"path": path},
        ) from exc
    if p.suffix.lower() != ".pptx":
        raise OfficeMCPError(
            ERR_UNSUPPORTED_FMT,
            f"Expected a .pptx file, got {p.suffix!r}",
            {"path": str(p), "extension": p.suffix},
        )
    return p


def _require_exists(path: Path) -> None:
    """Raise ``ERR_FILE_NOT_FOUND`` if ``path`` does not exist on disk."""
    if not path.exists():
        raise OfficeMCPError(
            ERR_FILE_NOT_FOUND,
            f"File not found: {path}",
            {"path": str(path)},
        )


def _open_pptx(path: Path) -> PresentationType:
    """Open an existing ``.pptx`` file.

    Combines :func:`_require_exists` with :class:`pptx.Presentation`
    so that callers always get a proper ``ERR_FILE_NOT_FOUND`` instead
    of a generic ``FileNotFoundError`` from python-pptx.
    """
    _require_exists(path)
    return Presentation(str(path))


def _save_pptx(prs: PresentationType, path: Path) -> None:
    """Save the presentation back to disk, creating parents if needed.

    Routes the underlying ``Presentation.save`` call through
    :func:`office_mcp.paths.save_with_lock_check` so a held lock
    (file open in PowerPoint) surfaces as ``ERR_FILE_LOCKED`` instead
    of an uncaught :class:`PermissionError`.
    """
    path.parent.mkdir(parents=True, exist_ok=True)
    save_with_lock_check(path, lambda: prs.save(str(path)))


def _check_slide_index(prs: PresentationType, index: int, *, op: str) -> None:
    """Validate a slide ``index`` against the deck's slide count.

    Raises :class:`OfficeMCPError` with ``ERR_INVALID_PARAMS`` when
    the index is not an int, is negative, or is out of range. ``op``
    is a short verb used in the error message ("read" / "delete" /
    etc.) so the message names the operation.
    """
    if not isinstance(index, int) or isinstance(index, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"slide index must be an int, got {type(index).__name__}",
            {"index": index},
        )
    n = len(prs.slides)
    if index < 0 or index >= n:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"slide index {index} is out of range for {op} "
                f"(deck has {n} slide(s); valid indices 0..{n - 1})"
            ),
            {"index": index, "count": n},
        )


def _check_layout_index(prs: PresentationType, layout_index: int) -> None:
    """Validate ``layout_index`` against the master's slide-layout count."""
    if not isinstance(layout_index, int) or isinstance(layout_index, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"layout_index must be an int, got {type(layout_index).__name__}",
            {"layout_index": layout_index},
        )
    layouts = prs.slide_layouts
    n = len(layouts)
    if layout_index < 0 or layout_index >= n:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"layout_index {layout_index} is out of range "
                f"(deck has {n} layout(s); valid indices 0..{n - 1})"
            ),
            {"layout_index": layout_index, "count": n},
        )


def _slide_title_text(slide: Slide) -> str:
    """Return the title placeholder's text, or ``""`` if none.

    Some layouts (e.g. ``"Blank"``) have no title placeholder;
    :pyattr:`slide.shapes.title` returns ``None`` for them. We treat
    that as an empty title rather than raising.
    """
    title = slide.shapes.title
    if title is None:
        return ""
    return title.text or ""


def _slide_shape_dict(shape) -> dict[str, Any]:
    """Return a JSON-friendly dict describing a single shape.

    ``type`` is the shape's class name (e.g. ``"TextBox"`` /
    ``"SlidePlaceholder"`` / ``"Picture"`` / ``"AutoShape"``) so the
    agent can distinguish text boxes from placeholders without
    importing python-pptx. ``name`` is the auto-generated PowerPoint
    name (``"TextBox 1"``, ``"Title 1"``, etc.). ``text`` is the
    text-frame text, or ``""`` if the shape has none.
    """
    text = ""
    try:
        if shape.has_text_frame:
            text = shape.text_frame.text or ""
    except (AttributeError, ValueError):
        text = ""
    return {
        "type": type(shape).__name__,
        "shape_type": int(shape.shape_type) if shape.shape_type is not None else None,
        "name": shape.name,
        "text": text,
    }


# ---------------------------------------------------------------------------
# 1. pptx_create_presentation
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_create_presentation(
    path: str,
    title: str | None = None,
    folder: str | None = None,
) -> dict[str, str]:
    """Create a new ``.pptx`` file at ``path``, optionally with a title slide.

    Args:
        path: Target ``.pptx`` path. May be absolute or relative to
            ``folder`` (or to the default folder when ``folder`` is
            ``None``).
        title: Optional text for the first slide's title placeholder.
            When ``None`` (the default), the file is created with no
            slides at all (``VAL-PPTX-002``). When a non-empty
            ``title`` is given, a single slide is appended using the
            ``"Title Slide"`` layout (index 0) and the title is
            written into its title placeholder (``VAL-PPTX-001``).
        folder: Optional base folder for relative paths.

    Returns:
        ``{"path": "<absolute path>"}``.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` if ``path`` is empty
            or a file already exists at the target,
            ``ERR_UNSUPPORTED_FMT`` if ``path`` does not end in
            ``.pptx``.
    """
    out = _resolve_pptx(path, folder)
    if out.exists():
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"File already exists: {out}",
            {"path": str(out)},
        )
    prs = Presentation()
    if title is not None:
        # Use layout 0 ("Title Slide") so the title placeholder is
        # the slide's only content placeholder. python-pptx
        # automatically gives us the title placeholder handle via
        # ``slide.shapes.title``.
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
    _save_pptx(prs, out)
    return {"path": str(out)}


# ---------------------------------------------------------------------------
# 2. pptx_get_info
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_get_info(
    path: str,
    folder: str | None = None,
) -> dict[str, Any]:
    """Return a summary of the presentation's structure and dimensions.

    The returned dict contains:

    * ``path`` — absolute path of the file on disk.
    * ``slide_count`` — number of slides (0 is valid;
      ``VAL-PPTX-009``).
    * ``layouts`` — list of layout names available in the deck's
      slide master (always non-empty; ``VAL-PPTX-006`` / ``-009``).
    * ``dimensions`` — ``{"width_inches": <float>, "height_inches":
      <float>}``. The default widescreen template is 13.333 × 7.5
      inches; the default ``Presentation()`` template is 10.0 × 7.5
      inches.

    Args:
        path: Path to an existing ``.pptx``.
        folder: Optional base folder for relative paths.

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if the file is
            missing, ``ERR_UNSUPPORTED_FMT`` for non-``.pptx``
            extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    layouts = [layout.name for layout in prs.slide_layouts]
    width_in = float(Emu(prs.slide_width).inches)
    height_in = float(Emu(prs.slide_height).inches)
    return {
        "path": str(p),
        "slide_count": len(prs.slides),
        "layouts": layouts,
        "dimensions": {
            "width_inches": width_in,
            "height_inches": height_in,
        },
    }


# ---------------------------------------------------------------------------
# 3. pptx_list_slides
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_list_slides(
    path: str,
    folder: str | None = None,
) -> list[dict[str, Any]]:
    """Return one dict per slide in insertion order.

    Each entry has the following keys (``VAL-PPTX-010``):

    * ``index`` — 0-based slide index.
    * ``layout`` — name of the slide's layout (e.g. ``"Title Slide"``,
      ``"Title and Content"``).
    * ``title`` — text of the slide's title placeholder, or ``""``
      when the layout has no title placeholder.
    * ``shape_count`` — number of shapes on the slide (placeholders
      + pictures + text boxes + ...).

    An empty deck returns ``[]`` (``VAL-PPTX-011``).

    Args:
        path: Path to an existing ``.pptx``.
        folder: Optional base folder for relative paths.

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if the file is
            missing, ``ERR_UNSUPPORTED_FMT`` for non-``.pptx``
            extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    out: list[dict[str, Any]] = []
    for idx, slide in enumerate(prs.slides):
        out.append(
            {
                "index": idx,
                "layout": slide.slide_layout.name,
                "title": _slide_title_text(slide),
                "shape_count": len(slide.shapes),
            }
        )
    return out


# ---------------------------------------------------------------------------
# 4. pptx_read_slide
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_read_slide(
    path: str,
    index: int,
    folder: str | None = None,
) -> dict[str, Any]:
    """Return the slide's title, layout, and shape list.

    Args:
        path: Path to an existing ``.pptx``.
        index: 0-based slide index.
        folder: Optional base folder for relative paths.

    Returns:
        ``{"index": <int>, "layout": <str>, "title": <str>,
        "shapes": [<shape_dict>, ...]}``.

        Each ``shape_dict`` has ``type`` (shape class name),
        ``shape_type`` (the numeric ``MSO_SHAPE_TYPE`` or ``None``),
        ``name`` (auto-generated name), and ``text`` (text-frame
        text, or ``""`` for shape types without text).

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for a non-int or
            out-of-range ``index``, ``ERR_FILE_NOT_FOUND`` if the
            file is missing, ``ERR_UNSUPPORTED_FMT`` for non-``.pptx``
            extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    _check_slide_index(prs, index, op="read")
    slide = prs.slides[index]
    return {
        "index": index,
        "layout": slide.slide_layout.name,
        "title": _slide_title_text(slide),
        "shapes": [_slide_shape_dict(shape) for shape in slide.shapes],
    }


# ---------------------------------------------------------------------------
# 5. pptx_add_slide
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_add_slide(
    path: str,
    layout_index: int = 1,
    title: str | None = None,
    folder: str | None = None,
) -> dict[str, int]:
    """Append a new slide to the deck.

    The slide is always appended to the end
    (``VAL-PPTX-018`` — the returned ``index`` equals the previous
    slide count). When ``title`` is given, the slide's title
    placeholder is populated (``VAL-PPTX-021``).

    Args:
        path: Path to an existing ``.pptx``.
        layout_index: 0-based index into ``prs.slide_layouts``. The
            default ``1`` is the ``"Title and Content"`` layout in
            python-pptx's default template.
        title: Optional title text. When ``None`` the title
            placeholder is left empty.
        folder: Optional base folder for relative paths.

    Returns:
        ``{"index": <int>}`` — the 0-based index of the new slide.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for a non-int or
            out-of-range ``layout_index`` (the deck is **not**
            modified in this case — ``VAL-PPTX-020``),
            ``ERR_FILE_NOT_FOUND`` if the file is missing,
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    _check_layout_index(prs, layout_index)
    layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(layout)
    if title is not None:
        # Some layouts (e.g. Blank) have no title placeholder; skip
        # silently in that case rather than crashing the call.
        if slide.shapes.title is not None:
            slide.shapes.title.text = title
    _save_pptx(prs, p)
    new_index = len(prs.slides) - 1
    return {"index": new_index}


# ---------------------------------------------------------------------------
# 6. pptx_delete_slide
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_delete_slide(
    path: str,
    index: int,
    folder: str | None = None,
) -> dict[str, bool]:
    """Remove a slide from the deck.

    Refuses to delete the last remaining slide (``VAL-PPTX-024``):
    the deck must always have at least one slide. After deletion
    the remaining slides keep their relative order, but indices
    past ``index`` shift down by one (``VAL-PPTX-023`` / ``-026``).

    Args:
        path: Path to an existing ``.pptx``.
        index: 0-based slide index.
        folder: Optional base folder for relative paths.

    Returns:
        ``{"ok": True}``.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for a non-int or
            out-of-range ``index``, or when the deck has only one
            slide; ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    n = len(prs.slides)
    if n <= 1:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"Cannot delete the only remaining slide (index {index}); "
                "a presentation must always have at least one slide."
            ),
            {"index": index, "slide_count": n},
        )
    _check_slide_index(prs, index, op="delete")
    # python-pptx does not expose a public delete-slide API, so we
    # remove the ``sldId`` element from the ``sldIdLst`` and drop
    # the relationship from the presentation part.
    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides_list = list(sldIdLst)
    target_sldId = slides_list[index]
    rId = target_sldId.get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    sldIdLst.remove(target_sldId)
    if rId:
        prs.part.drop_rel(rId)
    _save_pptx(prs, p)
    return {"ok": True}


# ---------------------------------------------------------------------------
# 7. pptx_reorder_slides
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_reorder_slides(
    path: str,
    from_index: int,
    to_index: int,
    folder: str | None = None,
) -> dict[str, bool]:
    """Move a slide from one position to another.

    This is a true *move*: the slide at ``from_index`` is removed
    first, then inserted at ``to_index`` in the resulting sequence
    (``VAL-PPTX-027``). So moving slide 0 to position 2 in a 3-slide
    deck ``[A, B, C]`` yields ``[B, C, A]``, not the swap result
    ``[C, B, A]``.

    Reordering a slide to its own index (``from_index == to_index``)
    is a no-op that still returns ``{"ok": True}``
    (``VAL-PPTX-028``).

    Args:
        path: Path to an existing ``.pptx``.
        from_index: 0-based index of the slide to move.
        to_index: 0-based destination index. Must satisfy
            ``0 <= to_index < slide_count`` (after the slide is
            removed from ``from_index``).
        folder: Optional base folder for relative paths.

    Returns:
        ``{"ok": True}``.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for a non-int or
            out-of-range ``from_index`` / ``to_index``, in which
            case the deck is left unchanged (``VAL-PPTX-029``);
            ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    n = len(prs.slides)
    if not isinstance(from_index, int) or isinstance(from_index, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"from_index must be an int, got {type(from_index).__name__}",
            {"from_index": from_index},
        )
    if not isinstance(to_index, int) or isinstance(to_index, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"to_index must be an int, got {type(to_index).__name__}",
            {"to_index": to_index},
        )
    if from_index < 0 or from_index >= n:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"from_index {from_index} is out of range "
                f"(deck has {n} slide(s); valid indices 0..{n - 1})"
            ),
            {"from_index": from_index, "count": n},
        )
    if to_index < 0 or to_index >= n:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"to_index {to_index} is out of range "
                f"(deck has {n} slide(s); valid indices 0..{n - 1})"
            ),
            {"to_index": to_index, "count": n},
        )
    if from_index == to_index:
        # No-op; do not rewrite the file so its SHA256 is preserved.
        return {"ok": True}

    sldIdLst = prs.slides._sldIdLst  # type: ignore[attr-defined]
    slides_list = list(sldIdLst)
    target_sldId = slides_list[from_index]
    # ``lxml`` semantics: ``sldIdLst.remove(elem)`` removes the
    # element from its parent, after which we re-insert it at the
    # requested new position. Python lists support ``.append()`` but
    # not ``.insert()`` for lxml elements — we use ``.addprevious()``
    # / ``.addnext()`` instead.
    sldIdLst.remove(target_sldId)
    # After removal, ``slides_list[from_index]`` no longer exists; we
    # need the element at ``to_index`` in the *current* sldIdLst.
    current_list = list(sldIdLst)
    if to_index == len(current_list):
        # Append at the end.
        sldIdLst.append(target_sldId)
    else:
        # Insert just before the slide currently at position
        # ``to_index`` (the desired destination).
        destination = current_list[to_index]
        destination.addprevious(target_sldId)
    _save_pptx(prs, p)
    return {"ok": True}


# ---------------------------------------------------------------------------
# Content / shape helpers
# ---------------------------------------------------------------------------


def _require_positive_number(name: str, value: float | int | None) -> float:
    """Validate that ``value`` is a finite number; return as ``float``.

    Used by the shape / text box / image geometry helpers below.
    """
    if value is None:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"{name} must be a number, got None",
            {name: value},
        )
    if isinstance(value, bool) or not isinstance(value, (int, float)):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"{name} must be a number, got {type(value).__name__}",
            {name: value, "type": type(value).__name__},
        )
    f = float(value)
    if f != f or f in (float("inf"), float("-inf")):  # NaN / Inf check
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"{name} must be a finite number, got {value!r}",
            {name: value},
        )
    return f


def _check_slide_for_edit(prs: PresentationType, slide: int) -> Slide:
    """Validate ``slide`` (0-based) and return the corresponding slide.

    Raises :class:`OfficeMCPError` (``ERR_INVALID_PARAMS``) for
    negative, non-int, or out-of-range indices; ``ERR_FILE_NOT_FOUND``
    is the caller's responsibility (this helper is only called after
    the file has been opened).
    """
    if not isinstance(slide, int) or isinstance(slide, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"slide must be an int, got {type(slide).__name__}",
            {"slide": slide, "type": type(slide).__name__},
        )
    n = len(prs.slides)
    if slide < 0 or slide >= n:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"slide index {slide} is out of range "
                f"(deck has {n} slide(s); valid indices 0..{n - 1})"
            ),
            {"slide": slide, "count": n},
        )
    return prs.slides[slide]


def _resolve_picture(image_path: str, folder: str | None) -> Path:
    """Resolve ``image_path`` (relative to ``folder`` if given) and check existence."""
    img_path = resolve_path(image_path, folder)
    if not img_path.exists():
        raise OfficeMCPError(
            ERR_FILE_NOT_FOUND,
            f"Image file not found: {img_path}",
            {"image_path": str(img_path)},
        )
    return img_path


def _picture_natural_size_emu(image_path: Path) -> tuple[int, int]:
    """Return the natural size of an image in EMU at 96 DPI.

    Uses :mod:`PIL.Image` to read the pixel dimensions, then
    multiplies by ``Inches(1/96)`` to convert to EMU. Returns
    ``(width_emu, height_emu)``.

    Falls back to a 1 inch square when the image cannot be read.
    """
    try:
        from PIL import Image  # type: ignore[import-not-found]
    except ImportError:
        # Pillow is a hard dep of the project; this branch is only
        # reached if someone strips it from the venv.
        return int(Inches(1.0)), int(Inches(1.0))
    try:
        with Image.open(str(image_path)) as img:
            px_w, px_h = img.size
    except (OSError, ValueError):
        return int(Inches(1.0)), int(Inches(1.0))
    width_emu = int(round(px_w / 96.0 * 914400))
    height_emu = int(round(px_h / 96.0 * 914400))
    return width_emu, height_emu


# ---------------------------------------------------------------------------
# 8. pptx_add_text_box
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_add_text_box(
    path: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    slide: int = 0,
    folder: str | None = None,
    font_size: float | None = None,
    bold: bool | None = None,
    italic: bool | None = None,
    font_name: str | None = None,
    color: str | None = None,
) -> dict[str, int]:
    """Add a text box with optional font properties to a slide.

    The text box is appended to the slide's shape list and its
    0-based index within that slide is returned (``VAL-PPTX-031``).
    When ``font_size`` / ``bold`` / ``italic`` / ``font_name`` /
    ``color`` are given, they are applied to the text box's first
    run (``VAL-PPTX-032``); when any of them is ``None`` the
    corresponding attribute is left at the layout / default value.

    Args:
        path: Path to an existing ``.pptx``.
        text: Text to put inside the text box.
        x: Left edge in inches (e.g. ``1.0``).
        y: Top edge in inches (e.g. ``2.0``).
        w: Width in inches (e.g. ``3.0``).
        h: Height in inches (e.g. ``0.5``).
        slide: 0-based slide index (``VAL-PPTX-034``).
        folder: Optional base folder for relative paths.
        font_size: Optional font size in points (``VAL-PPTX-032``).
        bold: Optional bold flag (``VAL-PPTX-032``).
        italic: Optional italic flag.
        font_name: Optional font name (e.g. ``"Arial"``).
        color: Optional hex color string (``"FF0000"``).

    Returns:
        ``{"shape_index": <int>}`` — 0-based index of the new
        text box within the slide's shape list.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for non-numeric or
            out-of-range geometry, out-of-range ``slide``, or invalid
            ``color``; ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    slide_obj = _check_slide_for_edit(prs, slide)
    xi = _require_positive_number("x", x)
    yi = _require_positive_number("y", y)
    wi = _require_positive_number("w", w)
    hi = _require_positive_number("h", h)

    tb = slide_obj.shapes.add_textbox(Inches(xi), Inches(yi), Inches(wi), Inches(hi))
    tf = tb.text_frame
    tf.text = text

    if (
        font_size is not None
        or bold is not None
        or italic is not None
        or font_name is not None
        or color is not None
    ):
        run = tf.paragraphs[0].runs[0] if tf.paragraphs[0].runs else None
        # ``add_textbox`` leaves the text frame with one paragraph
        # that contains one run; if the user-provided text is empty,
        # python-pptx may not create a run, so we make sure.
        if run is None:
            tf.paragraphs[0].add_run()
            run = tf.paragraphs[0].runs[0]
        if font_size is not None:
            run.font.size = Pt(float(font_size))
        if bold is not None:
            run.font.bold = bool(bold)
        if italic is not None:
            run.font.italic = bool(italic)
        if font_name is not None:
            run.font.name = font_name
        if color is not None:
            from pptx.dml.color import RGBColor  # local import

            try:
                run.font.color.rgb = RGBColor.from_string(color)
            except (ValueError, AttributeError) as exc:
                raise OfficeMCPError(
                    ERR_INVALID_PARAMS,
                    f"color must be a hex string like 'FF0000', got {color!r}",
                    {"color": color},
                ) from exc

    _save_pptx(prs, p)
    return {"shape_index": len(slide_obj.shapes) - 1}


# ---------------------------------------------------------------------------
# 9. pptx_add_image
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_add_image(
    path: str,
    image_path: str,
    x: float,
    y: float,
    w: float | None = None,
    h: float | None = None,
    slide: int = 0,
    folder: str | None = None,
) -> dict[str, int]:
    """Add an image to a slide at the given inch coordinates.

    The ``w`` / ``h`` arguments control the rendered size:

    * Both given — the picture is stretched to exactly those
      dimensions (``VAL-PPTX-036``).
    * Only ``w`` given — ``h`` is computed from the image's
      intrinsic aspect ratio (``VAL-PPTX-038``). Aspect ratio
      preserved within 1% tolerance.
    * Only ``h`` given — ``w`` is computed from the image's
      intrinsic aspect ratio (``VAL-PPTX-039``).
    * Neither given — the picture is rendered at its natural pixel
      size at 96 DPI (``VAL-PPTX-037``).

    Args:
        path: Path to an existing ``.pptx``.
        image_path: Path to the image to embed (any format python-pptx
            can read: PNG, JPEG, GIF, BMP, TIFF).
        x: Left edge in inches.
        y: Top edge in inches.
        w: Optional width in inches. ``None`` = compute from aspect
            ratio (when ``h`` is given) or use natural size.
        h: Optional height in inches. ``None`` = compute from aspect
            ratio (when ``w`` is given) or use natural size.
        slide: 0-based slide index.
        folder: Optional base folder for relative paths. Used for
            both the deck and the image path.

    Returns:
        ``{"shape_index": <int>}`` — 0-based index of the new picture
        within the slide's shape list.

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if either the deck
            or the image file is missing (``VAL-PPTX-040``);
            ``ERR_INVALID_PARAMS`` for out-of-range ``slide``
            (``VAL-PPTX-041``) or non-numeric geometry; ``ERR_UNSUPPORTED_FMT``
            for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    slide_obj = _check_slide_for_edit(prs, slide)
    xi = _require_positive_number("x", x)
    yi = _require_positive_number("y", y)

    img_path = _resolve_picture(image_path, folder)

    width_emu: int | None = None
    height_emu: int | None = None
    if w is not None and h is not None:
        wi = _require_positive_number("w", w)
        hi = _require_positive_number("h", h)
        width_emu = int(Inches(wi))
        height_emu = int(Inches(hi))
    elif w is not None:
        wi = _require_positive_number("w", w)
        width_emu = int(Inches(wi))
        # Aspect ratio from intrinsic pixel size
        nat_w, nat_h = _picture_natural_size_emu(img_path)
        if nat_w > 0 and nat_h > 0:
            height_emu = int(round(width_emu * nat_h / nat_w))
    elif h is not None:
        hi = _require_positive_number("h", h)
        height_emu = int(Inches(hi))
        nat_w, nat_h = _picture_natural_size_emu(img_path)
        if nat_w > 0 and nat_h > 0:
            width_emu = int(round(height_emu * nat_w / nat_h))
    else:
        # Both None — use natural pixel size at 96 DPI
        width_emu, height_emu = _picture_natural_size_emu(img_path)

    with open(img_path, "rb") as fh:
        image_bytes = io.BytesIO(fh.read())

    if width_emu is not None and height_emu is not None:
        slide_obj.shapes.add_picture(
            image_bytes,
            Inches(xi),
            Inches(yi),
            width=Emu(width_emu),
            height=Emu(height_emu),
        )
    elif width_emu is not None:
        slide_obj.shapes.add_picture(
            image_bytes,
            Inches(xi),
            Inches(yi),
            width=Emu(width_emu),
        )
    elif height_emu is not None:
        slide_obj.shapes.add_picture(
            image_bytes,
            Inches(xi),
            Inches(yi),
            height=Emu(height_emu),
        )
    else:  # pragma: no cover - _picture_natural_size_emu always returns a tuple
        slide_obj.shapes.add_picture(
            image_bytes, Inches(xi), Inches(yi)
        )

    _save_pptx(prs, p)
    return {"shape_index": len(slide_obj.shapes) - 1}


# ---------------------------------------------------------------------------
# 10. pptx_add_shape
# ---------------------------------------------------------------------------


#: Map of lowercase shape_type strings (e.g. ``"rectangle"``) to the
#: corresponding :class:`pptx.enum.shapes.MSO_SHAPE` enum member. The
#: mapping is case-insensitive on input but the keys are lowercase.
_SHAPE_TYPE_MAP: dict[str, MSO_SHAPE] = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "round_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "isoceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "isosceles_triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_triangle": MSO_SHAPE.RIGHT_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "parallelogram": MSO_SHAPE.PARALLELOGRAM,
    "trapezoid": MSO_SHAPE.TRAPEZOID,
    "hexagon": MSO_SHAPE.HEXAGON,
    "octagon": MSO_SHAPE.OCTAGON,
    "star_5_point": MSO_SHAPE.STAR_5_POINT,
    "star_6_point": MSO_SHAPE.STAR_6_POINT,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "left_right_arrow": MSO_SHAPE.LEFT_RIGHT_ARROW,
    "curved_right_arrow": MSO_SHAPE.CURVED_RIGHT_ARROW,
    "curved_left_arrow": MSO_SHAPE.CURVED_LEFT_ARROW,
    "curved_up_arrow": MSO_SHAPE.CURVED_UP_ARROW,
    "curved_down_arrow": MSO_SHAPE.CURVED_DOWN_ARROW,
    "callout_oval": MSO_SHAPE.OVAL_CALLOUT,
    "oval_callout": MSO_SHAPE.OVAL_CALLOUT,
    "callout_rectangle": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "rectangular_callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "heart": MSO_SHAPE.HEART,
    "lightning_bolt": MSO_SHAPE.LIGHTNING_BOLT,
    "sun": MSO_SHAPE.SUN,
    "moon": MSO_SHAPE.MOON,
    "cloud": MSO_SHAPE.CLOUD,
    "smiley_face": MSO_SHAPE.SMILEY_FACE,
}


def _resolve_shape_type(name: str) -> MSO_SHAPE:
    """Resolve a user-supplied shape type string to :class:`MSO_SHAPE`."""
    if not isinstance(name, str) or not name.strip():
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"shape_type must be a non-empty string, got {name!r}",
            {"shape_type": name},
        )
    key = name.strip().lower().replace("-", "_").replace(" ", "_")
    if key not in _SHAPE_TYPE_MAP:
        supported = sorted(_SHAPE_TYPE_MAP.keys())
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"Unknown shape_type {name!r}; supported types: "
                f"{', '.join(supported)}"
            ),
            {"shape_type": name, "supported": supported},
        )
    return _SHAPE_TYPE_MAP[key]


@mcp.tool()
def pptx_add_shape(
    path: str,
    shape_type: str,
    x: float,
    y: float,
    w: float,
    h: float,
    slide: int = 0,
    folder: str | None = None,
    text: str | None = None,
) -> dict[str, int]:
    """Add a shape (rectangle, oval, etc.) to a slide.

    ``shape_type`` is one of the case-insensitive names mapped by
    :data:`_SHAPE_TYPE_MAP` — the supported set covers the most
    common ``MSO_SHAPE`` members (``VAL-PPTX-042`` for rectangle,
    ``VAL-PPTX-043`` for oval and rounded rectangle). An unknown
    type returns ``ERR_INVALID_PARAMS`` and the deck is left
    unchanged (``VAL-PPTX-044``).

    Args:
        path: Path to an existing ``.pptx``.
        shape_type: Shape type (e.g. ``"rectangle"``, ``"oval"``,
            ``"rounded_rectangle"``). See :data:`_SHAPE_TYPE_MAP`.
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        slide: 0-based slide index (``VAL-PPTX-045``).
        folder: Optional base folder for relative paths.
        text: Optional text to put inside the shape's text frame.

    Returns:
        ``{"shape_index": <int>}`` — 0-based index of the new shape
        within the slide's shape list.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for an unknown shape
            type, out-of-range slide, or non-numeric geometry;
            ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    slide_obj = _check_slide_for_edit(prs, slide)
    mso_shape = _resolve_shape_type(shape_type)
    xi = _require_positive_number("x", x)
    yi = _require_positive_number("y", y)
    wi = _require_positive_number("w", w)
    hi = _require_positive_number("h", h)

    shape = slide_obj.shapes.add_shape(
        mso_shape, Inches(xi), Inches(yi), Inches(wi), Inches(hi)
    )
    if text is not None:
        shape.text_frame.text = text

    _save_pptx(prs, p)
    return {"shape_index": len(slide_obj.shapes) - 1}


# ---------------------------------------------------------------------------
# 11. pptx_add_table
# ---------------------------------------------------------------------------


def _validate_table_dims(rows: int, cols: int) -> None:
    """Validate ``rows`` and ``cols`` for :func:`pptx_add_table`."""
    if not isinstance(rows, int) or isinstance(rows, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"rows must be an int, got {type(rows).__name__}",
            {"rows": rows, "type": type(rows).__name__},
        )
    if not isinstance(cols, int) or isinstance(cols, bool):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"cols must be an int, got {type(cols).__name__}",
            {"cols": cols, "type": type(cols).__name__},
        )
    if rows <= 0:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"rows must be a positive integer, got {rows}",
            {"rows": rows},
        )
    if cols <= 0:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"cols must be a positive integer, got {cols}",
            {"cols": cols},
        )


def _pad_table_data(
    data: list[list[Any]], rows: int, cols: int
) -> list[list[Any]]:
    """Pad ``data`` to ``rows x cols`` with empty strings.

    The input is taken to be row-major: ``data[r][c]`` is the value
    for row ``r`` column ``c``. Missing rows and missing trailing
    columns are padded with ``""`` (``VAL-PPTX-048``). Non-string
    values are coerced to ``str`` so the table cell setter doesn't
    trip on ``int`` / ``float`` / ``None``.
    """
    out: list[list[Any]] = []
    for r in range(rows):
        row_in = data[r] if r < len(data) else None
        row_out: list[Any] = []
        for c in range(cols):
            if row_in is None:
                row_out.append("")
            elif c < len(row_in):
                v = row_in[c]
                row_out.append("" if v is None else str(v))
            else:
                row_out.append("")
        out.append(row_out)
    return out


@mcp.tool()
def pptx_add_table(
    path: str,
    rows: int,
    cols: int,
    x: float,
    y: float,
    w: float,
    h: float,
    slide: int = 0,
    folder: str | None = None,
    data: list[list[Any]] | None = None,
) -> dict[str, int]:
    """Add a table to a slide.

    When ``data`` is ``None`` the table is created with every cell
    empty (``VAL-PPTX-047``). When ``data`` is shorter than
    ``rows x cols`` the missing cells are padded with empty strings
    (``VAL-PPTX-048``). All cell values are coerced to ``str`` before
    being written.

    Args:
        path: Path to an existing ``.pptx``.
        rows: Number of rows (must be ``> 0``; ``VAL-PPTX-050``).
        cols: Number of columns (must be ``> 0``).
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        slide: 0-based slide index (``VAL-PPTX-049``).
        folder: Optional base folder for relative paths.
        data: Optional 2D list of cell values, row-major.

    Returns:
        ``{"shape_index": <int>}`` — 0-based index of the new table
        within the slide's shape list.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for non-positive
            ``rows`` / ``cols`` (``VAL-PPTX-050``), out-of-range
            ``slide`` (``VAL-PPTX-049``), or non-numeric geometry;
            ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    slide_obj = _check_slide_for_edit(prs, slide)
    _validate_table_dims(rows, cols)
    xi = _require_positive_number("x", x)
    yi = _require_positive_number("y", y)
    wi = _require_positive_number("w", w)
    hi = _require_positive_number("h", h)

    shape = slide_obj.shapes.add_table(
        rows, cols, Inches(xi), Inches(yi), Inches(wi), Inches(hi)
    )
    table = shape.table
    if data is not None:
        if not isinstance(data, list):
            raise OfficeMCPError(
                ERR_INVALID_PARAMS,
                f"data must be a list of lists, got {type(data).__name__}",
                {"data": data, "type": type(data).__name__},
            )
        padded = _pad_table_data(data, rows, cols)
        for r in range(rows):
            for c in range(cols):
                table.cell(r, c).text = padded[r][c]

    _save_pptx(prs, p)
    return {"shape_index": len(slide_obj.shapes) - 1}


# ---------------------------------------------------------------------------
# 12. pptx_add_chart
# ---------------------------------------------------------------------------


#: Map of lowercase chart_type strings to :class:`XL_CHART_TYPE` enum
#: members. The mapping is case-insensitive on input but the keys
#: are lowercase.
_CHART_TYPE_MAP: dict[str, XL_CHART_TYPE] = {
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
    "column_percent_stacked": XL_CHART_TYPE.COLUMN_STACKED_100,
    "bar": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
    "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
    "bar_percent_stacked": XL_CHART_TYPE.BAR_STACKED_100,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "line_stacked": XL_CHART_TYPE.LINE_STACKED,
    "pie": XL_CHART_TYPE.PIE,
    "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "area_stacked": XL_CHART_TYPE.AREA_STACKED,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
    "scatter_lines": XL_CHART_TYPE.XY_SCATTER_LINES,
    "scatter_smooth": XL_CHART_TYPE.XY_SCATTER_SMOOTH,
    "radar": XL_CHART_TYPE.RADAR,
    "radar_markers": XL_CHART_TYPE.RADAR_MARKERS,
    "radar_filled": XL_CHART_TYPE.RADAR_FILLED,
    "bubble": XL_CHART_TYPE.BUBBLE,
    "bubble_3d": XL_CHART_TYPE.BUBBLE_THREE_D_EFFECT,
}


def _resolve_chart_type(name: str) -> XL_CHART_TYPE:
    """Resolve a user-supplied chart type string to :class:`XL_CHART_TYPE`."""
    if not isinstance(name, str) or not name.strip():
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"chart_type must be a non-empty string, got {name!r}",
            {"chart_type": name},
        )
    key = name.strip().lower().replace("-", "_").replace(" ", "_")
    if key not in _CHART_TYPE_MAP:
        supported = sorted(_CHART_TYPE_MAP.keys())
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            (
                f"Unknown chart_type {name!r}; supported types: "
                f"{', '.join(supported)}"
            ),
            {"chart_type": name, "supported": supported},
        )
    return _CHART_TYPE_MAP[key]


def _build_chart_data(data: Any) -> CategoryChartData:
    """Convert a user-supplied chart data structure into :class:`CategoryChartData`.

    Expected input shape (per ``VAL-PPTX-051``):

    .. code-block:: python

       {
           "categories": ["Q1", "Q2", ...],
           "series": [
               {"name": "Revenue", "values": [10, 20, ...]},
               ...
           ],
       }

    Raises :class:`OfficeMCPError` with ``ERR_INVALID_PARAMS`` for
    missing keys, empty ``series`` list (``VAL-PPTX-054``), or any
    other structural problem.
    """
    if not isinstance(data, dict):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            f"data must be a dict, got {type(data).__name__}",
            {"data": data, "type": type(data).__name__},
        )
    categories = data.get("categories")
    if not isinstance(categories, list):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            "data.categories must be a list of strings",
            {"categories": categories},
        )
    series_in = data.get("series")
    if not isinstance(series_in, list):
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            "data.series must be a list of {name, values} dicts",
            {"series": series_in},
        )
    if len(series_in) == 0:
        raise OfficeMCPError(
            ERR_INVALID_PARAMS,
            "data.series must contain at least one series",
            {"series": series_in},
        )
    chart_data = CategoryChartData()
    chart_data.categories = [str(c) for c in categories]
    for idx, s in enumerate(series_in):
        if not isinstance(s, dict):
            raise OfficeMCPError(
                ERR_INVALID_PARAMS,
                f"series[{idx}] must be a dict, got {type(s).__name__}",
                {"series_index": idx, "series": s},
            )
        name = s.get("name")
        values = s.get("values")
        if not isinstance(name, str) or not name:
            raise OfficeMCPError(
                ERR_INVALID_PARAMS,
                f"series[{idx}].name must be a non-empty string",
                {"series_index": idx, "name": name},
            )
        if not isinstance(values, list):
            raise OfficeMCPError(
                ERR_INVALID_PARAMS,
                f"series[{idx}].values must be a list",
                {"series_index": idx, "values": values},
            )
        chart_data.add_series(name, list(values))
    return chart_data


@mcp.tool()
def pptx_add_chart(
    path: str,
    chart_type: str,
    data: dict[str, Any],
    x: float,
    y: float,
    w: float,
    h: float,
    slide: int = 0,
    folder: str | None = None,
) -> dict[str, int]:
    """Add a chart to a slide.

    The ``data`` argument must follow the structure documented in
    :func:`_build_chart_data`: a dict with ``"categories"`` (list of
    x-axis labels) and ``"series"`` (list of ``{"name", "values"}``
    dicts).

    Args:
        path: Path to an existing ``.pptx``.
        chart_type: One of ``column``, ``bar``, ``line``, ``pie``,
            etc. — see :data:`_CHART_TYPE_MAP`. (``VAL-PPTX-051``,
            ``VAL-PPTX-052``.)
        data: Chart data structure with ``categories`` and ``series``
            keys (``VAL-PPTX-051``).
        x: Left edge in inches.
        y: Top edge in inches.
        w: Width in inches.
        h: Height in inches.
        slide: 0-based slide index (``VAL-PPTX-055``).
        folder: Optional base folder for relative paths.

    Returns:
        ``{"shape_index": <int>}`` — 0-based index of the new chart
        within the slide's shape list.

    Raises:
        OfficeMCPError: ``ERR_INVALID_PARAMS`` for unknown
            ``chart_type`` (``VAL-PPTX-053``), empty ``series``
            (``VAL-PPTX-054``), out-of-range ``slide``
            (``VAL-PPTX-055``), or non-numeric geometry;
            ``ERR_FILE_NOT_FOUND`` if the file is missing;
            ``ERR_UNSUPPORTED_FMT`` for non-``.pptx`` extensions.
    """
    p = _resolve_pptx(path, folder)
    prs = _open_pptx(p)
    slide_obj = _check_slide_for_edit(prs, slide)
    xl_type = _resolve_chart_type(chart_type)
    chart_data = _build_chart_data(data)
    xi = _require_positive_number("x", x)
    yi = _require_positive_number("y", y)
    wi = _require_positive_number("w", w)
    hi = _require_positive_number("h", h)

    slide_obj.shapes.add_chart(
        xl_type,
        Inches(xi),
        Inches(yi),
        Inches(wi),
        Inches(hi),
        chart_data,
    )

    _save_pptx(prs, p)
    return {"shape_index": len(slide_obj.shapes) - 1}


# ---------------------------------------------------------------------------
# 13. pptx_export_pdf
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_export_pdf(
    path: str,
    output: str,
    folder: str | None = None,
) -> dict[str, str]:
    """Convert a ``.pptx`` file to PDF via LibreOffice headless.

    Delegates to :func:`office_mcp.exporters.export_to_pdf`. A unique
    ``-env:UserInstallation`` is used per call so multiple exports
    can run concurrently (``VAL-WORD-079`` pattern).

    Args:
        path: Path to an existing ``.pptx``.
        output: Target path for the produced PDF. The parent
            directory is created if it does not exist (``VAL-PPTX-057``).
        folder: Optional base folder for relative paths.

    Returns:
        ``{"output_path": "<absolute path of the produced PDF>"}``.

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if the source is
            missing (``VAL-PPTX-059``); ``ERR_UNSUPPORTED_FMT`` for
            non-``.pptx`` sources (``VAL-PPTX-060``);
            ``ERR_LIBREOFFICE_MISSING`` when ``soffice`` is not on
            PATH (``VAL-PPTX-058``); ``ERR_EXPORT_FAILED`` for any
            other failure.
    """
    from office_mcp import exporters

    src = _resolve_pptx(path, folder)
    _require_exists(src)
    out_path = resolve_path(output, folder)
    out_dir = out_path.parent
    produced = exporters.export_to_pdf(src, out_dir)
    if produced != out_path:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        produced.replace(out_path)
        final = out_path
    else:
        final = produced
    return {"output_path": str(final)}


# ---------------------------------------------------------------------------
# 14. pptx_export_html
# ---------------------------------------------------------------------------


@mcp.tool()
def pptx_export_html(
    path: str,
    output: str,
    folder: str | None = None,
) -> dict[str, str]:
    """Convert a ``.pptx`` file to HTML via LibreOffice headless.

    Delegates to :func:`office_mcp.exporters.export_to_html`. Unlike
    ``.docx`` (which uses mammoth), ``.pptx`` always requires
    LibreOffice (``VAL-PPTX-062``).

    Args:
        path: Path to an existing ``.pptx``.
        output: Target path for the produced HTML. The parent
            directory is created if it does not exist.
        folder: Optional base folder for relative paths.

    Returns:
        ``{"output_path": "<absolute path of the produced HTML>"}``.

    Raises:
        OfficeMCPError: ``ERR_FILE_NOT_FOUND`` if the source is
            missing (``VAL-PPTX-063``); ``ERR_UNSUPPORTED_FMT`` for
            non-``.pptx`` sources (``VAL-PPTX-064``);
            ``ERR_LIBREOFFICE_MISSING`` when ``soffice`` is not on
            PATH (``VAL-PPTX-062``); ``ERR_EXPORT_FAILED`` for any
            other failure.
    """
    from office_mcp import exporters

    src = _resolve_pptx(path, folder)
    _require_exists(src)
    out_path = resolve_path(output, folder)
    out_dir = out_path.parent
    produced = exporters.export_to_html(src, out_dir)
    if produced != out_path:
        out_path.parent.mkdir(parents=True, exist_ok=True)
        produced.replace(out_path)
        final = out_path
    else:
        final = produced
    return {"output_path": str(final)}
